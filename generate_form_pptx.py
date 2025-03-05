from pptx import Presentation
import qrcode
import os
from io import BytesIO

def generate_form_with_qr(template_path, output_path, qr_data):
    # Открываем шаблон презентации
    prs = Presentation(template_path)
    
    # Создаем QR-код
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=1,
    )
    qr.add_data(qr_data)
    qr.make(fit=True)
    qr_image = qr.make_image(fill_color="black", back_color="white")
    
    # Сохраняем QR-код во временный буфер
    image_stream = BytesIO()
    qr_image.save(image_stream, format='PNG')
    image_stream.seek(0)
    
    # Добавляем QR-код на первый слайд
    slide = prs.slides[0]
    
    # Ищем текст "МАРШРУТНАЯ КАРТА"
    target_text = "МАРШРУТНАЯ КАРТА"
    text_shape = None
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and target_text in shape.text:
            text_shape = shape
            break
    
    if text_shape:
        # Задаем размеры QR-кода
        qr_width = 400000  # ~0.4 см
        qr_height = 400000
        
        # Располагаем QR-код справа от текста на том же уровне
        left = text_shape.left + text_shape.width + 200000  # отступ от текста ~0.2 см
        # Выравниваем по вертикали с текстом
        top = text_shape.top + (text_shape.height - qr_height) / 2
        
        # Добавляем QR-код на слайд
        slide.shapes.add_picture(
            image_stream,
            left,
            top,
            width=qr_width,
            height=qr_height
        )
    else:
        # Если текст не найден, используем позицию по умолчанию
        slide_width = prs.slide_width
        qr_width = 400000
        qr_height = 400000
        left = slide_width - qr_width - 200000
        top = 100000
        slide.shapes.add_picture(
            image_stream,
            left,
            top,
            width=qr_width,
            height=qr_height
        )
    
    # Сохраняем результат
    prs.save(output_path)

def main():
    template_path = "ШАБЛОН.pptx"
    
    while True:
        print("\nГенератор QR-кодов для презентаций")
        print("-" * 40)
        qr_data = input("Введите данные для QR-кода (или 'выход' для завершения): ")
        
        if qr_data.lower() in ['выход', 'exit', 'quit', 'q']:
            print("Программа завершена.")
            break
        
        # Формируем имя выходного файла
        output_path = f"презентация_{len(qr_data)}.pptx"
        counter = 1
        while os.path.exists(output_path):
            output_path = f"презентация_{len(qr_data)}_{counter}.pptx"
            counter += 1
        
        try:
            generate_form_with_qr(template_path, output_path, qr_data)
            print(f"Файл успешно создан: {output_path}")
        except Exception as e:
            print(f"Произошла ошибка: {str(e)}")

if __name__ == "__main__":
    main() 