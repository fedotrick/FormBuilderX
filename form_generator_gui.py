from PySide6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, 
                             QHBoxLayout, QLabel, QLineEdit, QPushButton, 
                             QFileDialog, QFormLayout, QMessageBox, QGroupBox,
                             QDateEdit, QTimeEdit, QComboBox)
from PySide6.QtCore import Qt, QDate, QTime, QPropertyAnimation, QEasingCurve
from PySide6.QtGui import QFont
from pptx import Presentation
from pptx.util import Pt
import qrcode
import os
from io import BytesIO
import sys
from datetime import datetime
import subprocess  # Добавляем в начало файла
import time
import sqlite3
from create_history_db import save_form_data, validate_cluster_number, get_next_cluster_number

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Генератор маршрутных карт")
        self.setMinimumWidth(800)
        self.setMinimumHeight(900)
        
        # Устанавливаем общий стиль приложения
        self.setStyleSheet("""
            QMainWindow {
                background-color: #f5f5f5;
            }
            QGroupBox {
                font-family: 'Segoe UI', Arial;
                font-size: 14px;
                font-weight: bold;
                border: 2px solid #dcdcdc;
                border-radius: 6px;
                margin-top: 12px;
                padding-top: 10px;
            }
            QGroupBox::title {
                subcontrol-origin: margin;
                subcontrol-position: top center;
                padding: 0 10px;
                color: #2c3e50;
                background-color: #f5f5f5;
            }
            QLineEdit {
                padding: 8px;
                border: 1px solid #bdc3c7;
                border-radius: 4px;
                background-color: white;
                font-family: 'Segoe UI', Arial;
                font-size: 12px;
                min-height: 20px;
            }
            QLineEdit:focus {
                border: 2px solid #3498db;
            }
            QDateEdit, QTimeEdit {
                padding: 8px;
                border: 1px solid #bdc3c7;
                border-radius: 4px;
                background-color: white;
                font-family: 'Segoe UI', Arial;
                font-size: 12px;
                min-height: 20px;
            }
            QDateEdit:focus, QTimeEdit:focus {
                border: 2px solid #3498db;
            }
            QLabel {
                font-family: 'Segoe UI', Arial;
                font-size: 12px;
                color: #2c3e50;
            }
            QPushButton {
                background-color: #3498db;
                color: white;
                border: none;
                border-radius: 4px;
                padding: 8px 16px;
                font-family: 'Segoe UI', Arial;
                font-size: 13px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #2980b9;
            }
            QPushButton:pressed {
                background-color: #2472a4;
            }
        """)

        # Создаем центральный виджет и основной layout
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        main_layout = QVBoxLayout(central_widget)
        main_layout.setSpacing(20)
        main_layout.setContentsMargins(20, 20, 20, 20)

        # Подключаемся к базе данных
        self.db_conn = sqlite3.connect('справочник.db')
        self.db_cursor = self.db_conn.cursor()

        # Создаем группы полей
        self.create_cast_group()
        self.create_gluing_group()
        self.create_control_group()

        # Загружаем данные из справочника
        self.load_reference_data()

        # Добавляем группы в основной layout
        main_layout.addWidget(self.cast_group)
        main_layout.addWidget(self.gluing_group)
        main_layout.addWidget(self.control_group)

        # Добавляем кнопки
        button_layout = QHBoxLayout()
        button_layout.addStretch()
        self.generate_btn = QPushButton("Сгенерировать")
        self.generate_btn.setFixedWidth(200)
        self.generate_btn.setFixedHeight(40)
        self.generate_btn.clicked.connect(self.generate_form)
        button_layout.addWidget(self.generate_btn)
        main_layout.addLayout(button_layout)
        
        # Устанавливаем текущую дату и время
        self.set_current_datetime()

    def create_cast_group(self):
        self.cast_group = QGroupBox("Информация об отливке")
        layout = QFormLayout()
        layout.setSpacing(15)
        layout.setContentsMargins(20, 20, 20, 20)

        self.fields = {}
        
        # Создаем выпадающие списки для номера и наименования отливки
        self.fields['cast_number'] = QComboBox()
        self.fields['cast_name'] = QComboBox()
        self.fields['cluster_number'] = QLineEdit()

        # Добавляем подсказки
        self.fields['cast_number'].setPlaceholderText("Выберите номер отливки")
        self.fields['cast_name'].setPlaceholderText("Выберите наименование отливки")
        self.fields['cluster_number'].setPlaceholderText("Нажмите 'Сгенерировать номер'")

        # Связываем изменение номера с обновлением наименования
        self.fields['cast_number'].currentIndexChanged.connect(self.update_cast_name)

        # Создаем горизонтальный layout для поля номера кластера и кнопки
        number_layout = QHBoxLayout()
        number_layout.addWidget(self.fields['cluster_number'])
        
        # Создаем и добавляем кнопку генерации номера
        self.generate_number_btn = QPushButton("Сгенерировать номер")
        self.generate_number_btn.clicked.connect(self.generate_cluster_number)
        number_layout.addWidget(self.generate_number_btn)

        layout.addRow(self.create_label("Номер отливки (модели):"), self.fields['cast_number'])
        layout.addRow(self.create_label("Наименование отливки (модели):"), self.fields['cast_name'])
        layout.addRow(self.create_label("Номер кластера:"), number_layout)

        self.cast_group.setLayout(layout)

    def create_gluing_group(self):
        self.gluing_group = QGroupBox("Склейка элементов")
        layout = QFormLayout()
        layout.setSpacing(15)
        layout.setContentsMargins(20, 20, 20, 20)

        # Создаем поля для склейки
        self.fields['gluing_date'] = QDateEdit()
        self.fields['gluing_date'].setCalendarPopup(True)
        self.fields['gluing_date'].setDisplayFormat("dd.MM.yyyy")
        self.fields['gluing_executor'] = QComboBox()  # Меняем на QComboBox
        self.fields['gluing_quantity'] = QLineEdit()
        self.fields['gluing_notes'] = QLineEdit()

        # Добавляем подсказки
        self.fields['gluing_executor'].setPlaceholderText("Выберите исполнителя")
        self.fields['gluing_quantity'].setPlaceholderText("Введите количество")
        self.fields['gluing_notes'].setPlaceholderText("Введите примечания")

        layout.addRow(self.create_label("Дата:"), self.fields['gluing_date'])
        layout.addRow(self.create_label("Исполнитель:"), self.fields['gluing_executor'])
        layout.addRow(self.create_label("Количество:"), self.fields['gluing_quantity'])
        layout.addRow(self.create_label("Примечание:"), self.fields['gluing_notes'])

        self.gluing_group.setLayout(layout)

    def create_control_group(self):
        self.control_group = QGroupBox("Контроль сборки")
        layout = QFormLayout()
        layout.setSpacing(15)
        layout.setContentsMargins(20, 20, 20, 20)

        # Создаем поля для контроля
        self.fields['control_date'] = QDateEdit()
        self.fields['control_date'].setCalendarPopup(True)
        self.fields['control_date'].setDisplayFormat("dd.MM.yyyy")
        
        self.fields['control_time'] = QTimeEdit()
        self.fields['control_time'].setDisplayFormat("HH:mm")
        
        self.fields['control_executor'] = QComboBox()  # Меняем на QComboBox
        self.fields['control_quantity'] = QLineEdit()
        self.fields['control_notes'] = QLineEdit()

        # Добавляем подсказки
        self.fields['control_executor'].setPlaceholderText("Выберите контролера")
        self.fields['control_quantity'].setPlaceholderText("Введите количество")
        self.fields['control_notes'].setPlaceholderText("Введите примечания")

        layout.addRow(self.create_label("Дата:"), self.fields['control_date'])
        layout.addRow(self.create_label("Время:"), self.fields['control_time'])
        layout.addRow(self.create_label("Исполнитель:"), self.fields['control_executor'])
        layout.addRow(self.create_label("Количество:"), self.fields['control_quantity'])
        layout.addRow(self.create_label("Примечание:"), self.fields['control_notes'])

        self.control_group.setLayout(layout)

    def create_label(self, text):
        label = QLabel(text)
        # Убираем установку шрифта здесь, так как она теперь в общих стилях
        return label

    def set_current_datetime(self):
        current_date = QDate.currentDate()
        current_time = QTime.currentTime()
        
        self.fields['gluing_date'].setDate(current_date)
        self.fields['control_date'].setDate(current_date)
        self.fields['control_time'].setTime(current_time)

    def validate_fields(self):
        """Проверяет корректность заполнения полей"""
        try:
            # Проверяем обязательные поля
            required_fields = ['cluster_number']
            for field in required_fields:
                if not self.fields[field].text().strip():
                    return False, f"Поле '{field}' обязательно для заполнения"
            
            # Проверяем формат номера кластера
            validate_cluster_number(self.fields['cluster_number'].text().strip())
            
            return True, ""
        except ValueError as e:
            return False, str(e)

    def clear_fields(self):
        # Очищаем все поля, кроме дат и времени
        for field_name, widget in self.fields.items():
            if isinstance(widget, QLineEdit):
                widget.clear()
            elif isinstance(widget, (QDateEdit, QTimeEdit)):
                # Для полей даты и времени устанавливаем текущие значения
                continue
        
        # Обновляем текущую дату и время
        self.set_current_datetime()

    def generate_form(self):
        try:
            # Валидация полей
            is_valid, error_message = self.validate_fields()
            if not is_valid:
                QMessageBox.warning(self, "Предупреждение", error_message)
                return

            # Проверяем наличие файла шаблона
            if not os.path.exists("ШАБЛОН.pptx"):
                raise FileNotFoundError("Файл ШАБЛОН.pptx не найден в текущей директории")
            
            # Получаем данные из полей
            data = {}
            for field, widget in self.fields.items():
                if isinstance(widget, QDateEdit):
                    data[field] = widget.date().toString("dd.MM.yyyy")
                elif isinstance(widget, QTimeEdit):
                    data[field] = widget.time().toString("HH:mm")
                elif isinstance(widget, QComboBox):
                    data[field] = widget.currentText()
                else:
                    data[field] = widget.text()
            
            # Генерируем форму и получаем путь к файлу
            output_path = self.generate_pptx_with_data("ШАБЛОН.pptx", data)
            
            # Сохраняем данные в базу истории
            try:
                save_form_data(data)
            except ValueError as e:
                QMessageBox.warning(self, "Предупреждение", str(e))
                return
            
            # Печатаем файл
            try:
                if sys.platform == 'win32':  # Для Windows
                    # Используем команду для прямой печати файла
                    subprocess.run(['powershell', 'Start-Process', '-FilePath', output_path, 
                                  '-Verb', 'Print', '-WindowStyle', 'Hidden'], shell=True)
                    # Даем время на отправку на печать и закрываем PowerPoint
                    time.sleep(2)  # Ждем 2 секунды
                    try:
                        subprocess.run(['taskkill', '/F', '/IM', 'POWERPNT.EXE'], 
                                     shell=True, 
                                     stderr=subprocess.DEVNULL,
                                     stdout=subprocess.DEVNULL)
                    except:
                        pass  # Игнорируем ошибку, если PowerPoint уже закрыт
                else:  # Для Linux/Mac
                    subprocess.run(['lpr', output_path])
            except Exception as e:
                print(f"Ошибка при печати: {str(e)}")
                QMessageBox.warning(self, "Предупреждение", 
                                  "Файл создан, но не удалось отправить на печать автоматически.\n"
                                  f"Файл сохранен как: {output_path}")
            
            # Очищаем поля после успешной генерации и печати
            self.clear_fields()
            
            # Устанавливаем текущую дату и время после очистки
            self.set_current_datetime()
            
        except Exception as e:
            QMessageBox.critical(self, "Ошибка", f"Неожиданная ошибка: {str(e)}")

    def generate_pptx_with_data(self, template_path, data):
        # Открываем шаблон
        prs = Presentation(template_path)
        
        # Проверяем наличие слайдов
        if not prs.slides:
            raise ValueError("Презентация не содержит слайдов")
            
        slide = prs.slides[0]
        
        # Создаем QR-код с обработкой ошибок
        try:
            qr = qrcode.QRCode(
                version=1,
                error_correction=qrcode.constants.ERROR_CORRECT_L,
                box_size=10,
                border=1,
            )
            print(f"Создание QR-кода для значения: {data['cluster_number']}")
            qr.add_data(str(data['cluster_number']))  # Явно преобразуем в строку
            qr.make(fit=True)
            qr_image = qr.make_image(fill_color="black", back_color="white")
        except Exception as e:
            print(f"Ошибка при создании QR-кода: {str(e)}")
            raise

        # Сохраняем QR-код во временный буфер
        image_stream = BytesIO()
        qr_image.save(image_stream, format='PNG')
        image_stream.seek(0)

        # Добавляем QR-код рядом с "МАРШРУТНАЯ КАРТА"
        for shape in slide.shapes:
            if hasattr(shape, "text") and "МАРШРУТНАЯ КАРТА" in shape.text:
                qr_width = 400000
                qr_height = 400000
                left = shape.left + shape.width + 200000
                top = shape.top + (shape.height - qr_height) / 2
                slide.shapes.add_picture(
                    image_stream,
                    left,
                    top,
                    width=qr_width,
                    height=qr_height
                )

        # Работаем с таблицами
        tables_found = False
        for shape in slide.shapes:
            if shape.has_table:
                tables_found = True
                table = shape.table
                print(f"Найдена таблица: {len(table.rows)} строк, {len(table.columns)} столбцов")
                
                # Проверяем первую ячейку таблицы, чтобы определить тип таблицы
                table_header = table.cell(0, 0).text.strip()
                print(f"Заголовок таблицы: '{table_header}'")

                # Если это таблица с номером отливки
                if "Номер" in table_header:
                    print("Обрабатываем таблицу с номером отливки")
                    try:
                        # Заполняем данные отливки с маленьким жирным шрифтом
                        for col in range(3):
                            cell = table.cell(1, col)
                            paragraph = cell.text_frame.paragraphs[0]
                            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                            if col == 0:
                                run.text = data['cast_number']
                            elif col == 1:
                                run.text = data['cast_name']
                            else:
                                run.text = data['cluster_number']
                            run.font.size = Pt(9)
                            run.font.bold = True  # Делаем шрифт жирным
                    except Exception as e:
                        print(f"Ошибка при заполнении данных отливки: {str(e)}")
                    continue  # Переходим к следующей таблице

                # Если это таблица операций
                try:
                    # Сначала найдем индексы нужных столбцов
                    column_indices = {}
                    print("\nЗаголовки столбцов:")
                    for col in range(len(table.columns)):
                        header = table.cell(0, col).text.strip()
                        print(f"Столбец {col}: '{header}'")
                        if header == "Дата":
                            column_indices['date'] = col
                        elif header == "Время":
                            column_indices['time'] = col
                        elif header == "Исполнитель":
                            column_indices['executor'] = col
                        elif header == "Количество":
                            column_indices['quantity'] = col
                        elif header == "Примечание":
                            column_indices['notes'] = col
                    
                    print("\nНайденные индексы столбцов:", column_indices)

                    # Теперь найдем строки операций и заполним данные
                    print("\nСодержимое первого столбца:")
                    for row in range(len(table.rows)):
                        try:
                            operation = table.cell(row, 0).text.strip()
                            print(f"Строка {row}: '{operation}'")
                            
                            # Заполняем данные для Склейки
                            if "Склейка элементов п/м" in operation:
                                print("Найдена строка Склейки")
                                for col_name, col_idx in column_indices.items():
                                    cell = table.cell(row, col_idx)
                                    paragraph = cell.text_frame.paragraphs[0]
                                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                                    if col_name == 'date':
                                        run.text = data['gluing_date']
                                    elif col_name == 'executor':
                                        run.text = data['gluing_executor']
                                    elif col_name == 'quantity':
                                        run.text = data['gluing_quantity']
                                    elif col_name == 'notes':
                                        run.text = data['gluing_notes']
                                    run.font.size = Pt(9)
                                    run.font.bold = True  # Делаем шрифт жирным
                            
                            # Заполняем данные для Контроля
                            if "Контроль сборки кластера" in operation:
                                print("Найдена строка Контроля")
                                for col_name, col_idx in column_indices.items():
                                    cell = table.cell(row, col_idx)
                                    paragraph = cell.text_frame.paragraphs[0]
                                    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                                    if col_name == 'date':
                                        run.text = data['control_date']
                                    elif col_name == 'executor':
                                        run.text = data['control_executor']
                                    elif col_name == 'quantity':
                                        run.text = data['control_quantity']
                                    elif col_name == 'notes':
                                        run.text = data['control_notes']
                                    run.font.size = Pt(9)
                                    run.font.bold = True  # Делаем шрифт жирным
                                
                                # Для времени контроля
                                for r in range(len(table.rows)):
                                    for c in range(len(table.columns)):
                                        if table.cell(r, c).text.strip() == "Время:":
                                            cell = table.cell(r, c + 1)
                                            paragraph = cell.text_frame.paragraphs[0]
                                            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                                            run.text = data['control_time']
                                            run.font.size = Pt(9)
                                            run.font.bold = True  # Делаем шрифт жирным

                        except Exception as e:
                            print(f"Ошибка при обработке строки {row}: {str(e)}")
                            
                except Exception as e:
                    print(f"Ошибка при обработке таблицы: {str(e)}")
        
        if not tables_found:
            print("На слайде не найдено ни одной таблицы!")

        # Формируем имя выходного файла, заменяя недопустимые символы
        cluster_number = data['cluster_number'].replace('/', '_').replace('\\', '_')
        output_path = f"маршрутная_карта_{cluster_number}.pptx"
        counter = 1
        while os.path.exists(output_path):
            output_path = f"маршрутная_карта_{cluster_number}_{counter}.pptx"
            counter += 1

        # Сохраняем результат
        prs.save(output_path)
        
        # Возвращаем путь к сохраненному файлу
        return output_path

    def show(self):
        self.setWindowOpacity(1.0)  # Устанавливаем непрозрачность сразу
        super().show()  # Показываем окно

    def load_reference_data(self):
        """Загрузка справочных данных из базы"""
        try:
            # Загружаем данные отливок
            self.db_cursor.execute('SELECT "Номер", "Наименование" FROM "ЛГМ_Отливки"')
            lgm_casts = self.db_cursor.fetchall()
            self.db_cursor.execute('SELECT "Номер", "Наименование" FROM "ЛПД_Отливки"')
            lpd_casts = self.db_cursor.fetchall()
            self.db_cursor.execute('SELECT "Наименование" FROM "Прочие_Отливки"')
            other_casts = self.db_cursor.fetchall()

            # Заполняем выпадающие списки номеров и наименований отливок
            self.cast_numbers_data = {}  # Словарь для хранения соответствия номер-наименование
            
            # Очищаем списки перед заполнением
            self.fields['cast_number'].clear()
            self.fields['cast_name'].clear()
            
            # Добавляем ЛГМ и ЛПД отливки
            for number, name in lgm_casts + lpd_casts:
                self.fields['cast_number'].addItem(number)
                self.fields['cast_name'].addItem(name)
                self.cast_numbers_data[number] = name
            
            # Добавляем прочие отливки
            for (name,) in other_casts:
                self.fields['cast_number'].addItem(name)
                self.fields['cast_name'].addItem(name)
                self.cast_numbers_data[name] = name

            # Загружаем сборщиков
            self.db_cursor.execute('SELECT "ФИО" FROM "Сборщики"')
            assemblers = [row[0] for row in self.db_cursor.fetchall()]
            self.fields['gluing_executor'].addItems(assemblers)

            # Загружаем контролеров сборки
            self.db_cursor.execute('SELECT "ФИО" FROM "Контролеры_Сборки"')
            controllers = [row[0] for row in self.db_cursor.fetchall()]
            self.fields['control_executor'].addItems(controllers)

        except Exception as e:
            print(f"Ошибка при загрузке справочных данных: {e}")

    def update_cast_name(self, index):
        """Обновляет поле наименования при выборе номера отливки"""
        current_number = self.fields['cast_number'].currentText()
        if current_number in self.cast_numbers_data:
            name = self.cast_numbers_data[current_number]
            index = self.fields['cast_name'].findText(name)
            if index >= 0:
                self.fields['cast_name'].setCurrentIndex(index)

    def generate_cluster_number(self):
        """Генерирует и устанавливает номер кластера на основе даты склейки"""
        try:
            # Получаем дату склейки
            gluing_date = self.fields['gluing_date'].date().toString("dd.MM.yyyy")
            
            # Генерируем следующий номер кластера
            next_number = get_next_cluster_number(gluing_date)
            
            # Устанавливаем номер в поле
            self.fields['cluster_number'].setText(next_number)
            
        except Exception as e:
            QMessageBox.warning(self, "Предупреждение", str(e))

    def closeEvent(self, event):
        """Закрываем соединение с базой при закрытии приложения"""
        try:
            if hasattr(self, 'db_conn'):
                self.db_conn.close()
        except Exception as e:
            print(f"Ошибка при закрытии соединения с БД: {e}")
        event.accept()

def main():
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec())

if __name__ == "__main__":
    main() 